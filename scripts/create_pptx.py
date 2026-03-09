#!/usr/bin/env python3
"""
Create a prototype-snapshots PPTX slide in the style of the reference image.
Dark background, 5 screenshots in a collage layout with styled labels.
"""

import os
import math
from io import BytesIO

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────────────────────────────────────
SCREENSHOTS_DIR = "/Users/nidhimaru/Developer/AI_for_Bharat-Kiro-submission/temp/screenshots"
OUTPUT_PATH = "/Users/nidhimaru/Developer/AI_for_Bharat-Kiro-submission/temp/prototype_snapshots.pptx"

# Slide size: 16:9 widescreen (13.333 x 7.5 inches)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
BG_COLOR       = RGBColor(0x03, 0x07, 0x12)   # #030712 void black
LABEL_YELLOW   = RGBColor(0xFF, 0xE5, 0x00)   # bright yellow
LABEL_LIME     = RGBColor(0xA3, 0xE6, 0x35)   # lime green
LABEL_CYAN     = RGBColor(0x22, 0xD3, 0xEE)   # cyan
LABEL_ORANGE   = RGBColor(0xF9, 0x73, 0x16)   # orange
LABEL_PURPLE   = RGBColor(0x81, 0x8C, 0xF8)   # indigo/purple
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT      = RGBColor(0x0A, 0x0E, 0x1A)   # near black for dark text on light bg


# ─────────────────────────────────────────────────────────────────────────────
# HELPER: add cropped/rounded screenshot with shadow to slide
# ─────────────────────────────────────────────────────────────────────────────
def add_screenshot(slide, img_path: str, left: float, top: float,
                   width: float, height: float, rotation: float = 0.0,
                   border_color: RGBColor | None = None):
    """
    Add a screenshot image to the slide.
    Applies a subtle inner glow border and saves as PNG into a BytesIO stream.
    left/top/width/height are in inches.
    """
    img = Image.open(img_path).convert("RGBA")

    # Crop to match target aspect ratio
    target_ratio = width / height
    src_w, src_h = img.size
    src_ratio = src_w / src_h

    if src_ratio > target_ratio:
        # Image is wider than target – crop sides
        new_w = int(src_h * target_ratio)
        offset_x = (src_w - new_w) // 2
        img = img.crop((offset_x, 0, offset_x + new_w, src_h))
    elif src_ratio < target_ratio:
        # Image is taller than target – crop bottom
        new_h = int(src_w / target_ratio)
        img = img.crop((0, 0, src_w, new_h))

    # Resize for quality (at 2x for retina, will be down-sampled by PowerPoint)
    target_px_w = int(width * 150)   # 150 dpi rendering
    target_px_h = int(height * 150)
    img = img.resize((target_px_w, target_px_h), Image.LANCZOS)

    # Add rounded corners
    radius = max(8, int(target_px_w * 0.012))
    mask = Image.new("L", img.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle([(0, 0), (img.width - 1, img.height - 1)], radius=radius, fill=255)
    img.putalpha(mask)

    # Add glowing border around rounded rect
    if border_color:
        border_layer = Image.new("RGBA", img.size, (0, 0, 0, 0))
        bd = ImageDraw.Draw(border_layer)
        bc = (border_color[0], border_color[1], border_color[2], 180)
        bd.rounded_rectangle(
            [(1, 1), (img.width - 2, img.height - 2)],
            radius=radius, outline=bc, width=2
        )
        img = Image.alpha_composite(img, border_layer)

    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)

    pic = slide.shapes.add_picture(
        buf,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    if rotation != 0:
        pic.rotation = rotation

    return pic


def add_label(slide, text: str, left: float, top: float, width: float,
              bg_color: RGBColor, text_color: RGBColor = None,
              font_size: int = 14, rotation: float = 0.0):
    """
    Add a bold badge-style label text box.
    left/top/width are in inches. Height is auto-fitted.
    """
    if text_color is None:
        # Auto-pick dark or light text based on bg brightness
        r, g, b = bg_color[0], bg_color[1], bg_color[2]
        brightness = (r * 299 + g * 587 + b * 114) / 1000
        text_color = DARK_TEXT if brightness > 128 else WHITE

    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(0.55)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color

    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.name = "Arial Black"

    if rotation != 0:
        box.rotation = rotation

    return box


def set_slide_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
def create_pptx():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BG_COLOR)

    ss = SCREENSHOTS_DIR

    # ──────────────────────────────────────────────────────────────────────────
    # LAYOUT: Reference-style collage
    #
    #   ┌────────────────────┐  ┌──────────────────────┐
    #   │                    │  │                       │
    #   │    DASHBOARD       │  │   CONTENT GENERATOR   │
    #   │   (large left)     │  │      (large right)    │
    #   │                    │  │                       │
    #   └────────────────────┘  └──────────────────────┘
    #   [LABEL yellow]           [LABEL lime]
    #
    #   [LABEL cyan]   [LABEL orange]
    #   ┌──────────┐  ┌──────────────┐
    #   │ CALENDAR │  │  ANALYTICS   │
    #   │ (tilted) │  │  (straight)  │
    #   └──────────┘  └──────────────┘
    # ──────────────────────────────────────────────────────────────────────────

    # ── Top-left: Dashboard ──────────────────────────────────────────────────
    add_screenshot(
        slide,
        img_path=os.path.join(ss, "dashboard.png"),
        left=0.22, top=0.25, width=6.3, height=3.95,
        rotation=0,
        border_color=RGBColor(0x63, 0x66, 0xF1),
    )
    add_label(
        slide,
        text="Creator Dashboard",
        left=0.22, top=4.28, width=6.3,
        bg_color=LABEL_YELLOW, text_color=DARK_TEXT,
        font_size=13,
    )

    # ── Top-right: Content Generator (results) ───────────────────────────────
    add_screenshot(
        slide,
        img_path=os.path.join(ss, "content_generator.png"),
        left=6.82, top=0.15, width=6.3, height=4.05,
        rotation=0,
        border_color=RGBColor(0xA3, 0xE6, 0x35),
    )
    add_label(
        slide,
        text="AI Content Generator",
        left=6.82, top=4.28, width=6.3,
        bg_color=LABEL_LIME, text_color=DARK_TEXT,
        font_size=13,
    )

    # ── Bottom-left: Calendar (slight tilt left) ─────────────────────────────
    add_label(
        slide,
        text="Content Calendar",
        left=0.22, top=4.82, width=4.6,
        bg_color=LABEL_CYAN, text_color=DARK_TEXT,
        font_size=11,
        rotation=-1.5,
    )
    add_screenshot(
        slide,
        img_path=os.path.join(ss, "calendar.png"),
        left=0.22, top=5.45, width=4.6, height=1.88,
        rotation=-1.5,
        border_color=RGBColor(0x22, 0xD3, 0xEE),
    )

    # ── Bottom-right: Analytics (slight tilt right) ───────────────────────────
    add_label(
        slide,
        text="Analytics & Insights",
        left=5.12, top=4.82, width=7.98,
        bg_color=LABEL_ORANGE, text_color=WHITE,
        font_size=11,
        rotation=1.5,
    )
    add_screenshot(
        slide,
        img_path=os.path.join(ss, "analytics.png"),
        left=5.12, top=5.45, width=7.98, height=1.88,
        rotation=1.5,
        border_color=RGBColor(0xF9, 0x73, 0x16),
    )

    # ── Save ──
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"PPTX saved to: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_pptx()
